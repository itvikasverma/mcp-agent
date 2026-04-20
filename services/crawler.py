import hashlib
import os
from pathlib import Path
import base64
import pdfplumber
import io
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import openpyxl
from typing import List, Dict
import requests
import pandas as pd
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    TextLoader, UnstructuredImageLoader
)
from services.llm_groq import llm


# ───── Config ───── #
API_URL = "http://127.0.0.1:9000/add_doc"
SOURCE_ID = "confluence"
BASE_DIR = "data/store"
FOLDER_TO_PROCESS = f"{BASE_DIR}/{SOURCE_ID}"
MAX_CHUNK_WORDS = 2000  # Words per chunk

# ───── Helper Functions ───── #
def generate_title_from_content(content: str) -> str:
    prompt = (
        "Please create a good title for the following document. "
        "Don't explain anything, just return the title:\n\n"
        f"{content[:2000]}"
    )
    response = llm.invoke(prompt)
    return getattr(response, "content", str(response)).strip()


def extract_text_from_csv(file_path: str) -> str:
    try:
        df = pd.read_csv(file_path)
    except Exception:
        df = pd.read_csv(file_path, encoding="ISO-8859-1")
    return "\n".join([" | ".join(map(str, row)) for row in df.values])

def call_llm(bytes):
    return f"[PDF Byes]: {bytes[:50]}..."
    

def chunk_text(text: str, max_words: int):
    words = text.split()
    for i in range(0, len(words), max_words):
        yield ' '.join(words[i:i + max_words])


def load_pdf(file_path):
    results = []

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            blocks = []
            page_width, page_height = page.width, page.height

            # --- Extract text boxes with coordinates ---
            for line in page.extract_words(x_tolerance=2, y_tolerance=2, keep_blank_chars=False):
                top = max(0, min(line["top"], page_height))
                blocks.append({
                    "type": "text",
                    "content": line["text"],
                    "top": top
                })

            # --- Extract images with coordinates ---
            for img in page.images:
                # Clip coordinates to page boundaries
                x0 = max(0, min(img['x0'], page_width))
                x1 = max(0, min(img['x1'], page_width))
                top = max(0, min(img['top'], page_height))
                bottom = max(0, min(img['bottom'], page_height))

                # Crop and convert to PNG
                cropped_image = page.crop((x0, top, x1, bottom)).to_image(resolution=150)
                buf = io.BytesIO()
                cropped_image.save(buf, format="PNG")
                image_bytes = buf.getvalue()
                image_b64 = base64.b64encode(image_bytes).decode("utf-8")

                # send to LLM for description
                image_text = call_llm(image_b64)

                blocks.append({
                    "type": "text", 
                    "content": image_text,
                    "top": top
                })

            # --- Sort all blocks by vertical position (top) ---
            blocks.sort(key=lambda b: b["top"])

            # --- Merge consecutive text blocks ---
            merged_blocks = []
            buffer = ""
            for block in blocks:
                if block["type"] == "text":
                    buffer += block["content"] + " "
                else:
                    if buffer:
                        merged_blocks.append({"type": "text", "content": buffer.strip(), "page": page_num})
                        buffer = ""
                    merged_blocks.append({"type": "text", "content": block["content"], "page": page_num})
            if buffer:
                merged_blocks.append({"type": "text", "content": buffer.strip(), "page": page_num})

            results.extend(merged_blocks)

    return results


def load_docx(path: str) -> List[Dict]:
    results = []
    doc = docx.Document(path)

    for para in doc.paragraphs:
        # Add paragraph text if present
        text = para.text.strip()
        if text:
            results.append({"type": "text", "content": text})

        # Check for inline images in runs
        for run in para.runs:
            if "graphic" in run.element.xml:
                # Get image rId
                rId = run.element.xpath(".//a:blip/@r:embed")  # remove namespaces argument
                if rId:
                    part = doc.part.related_parts[rId[0]]
                    image_bytes = part.blob

                    # Print first 4 bytes as a preview
                    preview = " ".join(f"{b:02x}" for b in image_bytes[:4])
                    print(f"[INFO] image bytes preview: {preview} ...")

                    # Convert to base64 for LLM input
                    image_b64 = base64.b64encode(image_bytes).decode("utf-8")
                    print(f"image bytes {image_b64[:50]}...")  # preview only


                    # meaning = llm(image_b64)

                    # Append LLM meaning as text
                    # results.append({"type": "text", "content": meaning})

    return results





def load_pptx(path: str) -> List[Dict]:
    results = []
    prs = Presentation(path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        for img_index, shape in enumerate(slide.shapes, start=1):
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    results.append({"type": "text", "content": text})

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # Get raw bytes
                    image_bytes = shape.image.blob

                    # Convert to base64
                    image_b64 = base64.b64encode(image_bytes).decode("utf-8")

                    # Print first 100 characters of base64 for preview
                    print(f"[PPTX IMAGE BASE64] Slide {slide_num}, image {img_index}: {image_b64[:100]}...")

                    # Optionally send to LLM
                    # meaning = llm(image_b64, is_image=True)
                    # results.append({"type": "text", "content": meaning})

                except Exception as e:
                    print(f"[PPTX IMAGE ERROR] Slide {slide_num}, image {img_index}: {e}")

    return results


def load_xlsx(path: str) -> List[Dict]:
    results = []
    wb = openpyxl.load_workbook(path, data_only=True)

    for sheet in wb.worksheets:
        # Go row by row
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) if cell else "" for cell in row])
            if row_text.strip():
                results.append({"type": "text", "content": row_text})

        # Handle images (Excel stores them separately, no inline order)
        if hasattr(sheet, "_images"):
            for img in sheet._images:
                try:
                    img_bytes = img._data()
                except AttributeError:
                    img_bytes = img.ref.blob
                    image_b64 = base64.b64encode(img_bytes).decode("utf-8")
        
                    print(f"[INFO] image bytes: {str(image_b64)[:4]} ...")
                # meaning = llm(img_bytes)
                # results.append({"type": "text", "content": meaning})

    return results

# ───── Generic Document Loader ───── #
def load_document(file_path: str):
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            return load_pdf(file_path)
        elif ext == ".csv":
            content = extract_text_from_csv(file_path)
            return [Document(page_content=content, metadata={"source": file_path})]
        elif ext in [".txt", ".md"]:
            return TextLoader(file_path).load()
        elif ext in [".doc", ".docx"]:
            return load_docx(file_path)
        elif ext in [".html", ".htm", ".pptx", ".epub"]:
            return load_pptx(file_path)
        elif ext in [".xlsx", ".xls"]:
            return load_xlsx(file_path)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif"]:
            return UnstructuredImageLoader(file_path).load()
        else:
            raise ValueError(f"Unsupported file type: {file_path}")
    except Exception as e:
        print(f"[Error loading] {file_path}: {e}")
        return []


# ───── Processing ───── #
def process_single_file(file_path: str):
    documents = load_document(file_path)
    if not documents:
        print(f"[SKIP] No content: {file_path}")
        return

    full_content = "\n\n".join(
        d["content"] for d in documents if d["type"] == "text" and d["content"].strip()
    )
    if not full_content:
        print(f"[SKIP] Empty content: {file_path}")
        return

    # -----------------------------
    # 1️⃣ Generate content-based SHA256 hash
    # -----------------------------
    file_hash = hashlib.sha256(full_content.encode("utf-8")).hexdigest()
    print("File Hash:", file_hash)

    # -----------------------------
    # 2️⃣ Check if this hash already exists in your database
    # -----------------------------
    try:
        res = requests.get(f"http://127.0.0.1:9000/hash_list?file_hash={file_hash}")
        if res.status_code == 200 and res.json().get("exists", False):
            print(f"[SKIP] Already chunked: {file_path}")
            return
    except Exception as e:
        print(f"[WARN] Could not check hash for {file_path}: {e}")

    # -----------------------------
    # 3️⃣ Generate file title
    # -----------------------------
    title = generate_title_from_content(full_content)

    # -----------------------------
    # 4️⃣ Chunk text
    # -----------------------------
    chunks = list(chunk_text(full_content, MAX_CHUNK_WORDS))

    for idx, chunk in enumerate(chunks, start=1):
        safe_chunk = chunk.encode('utf-8', errors='replace').decode('utf-8')

        payload = {
            "doc_title": title,
            "content": safe_chunk,
            "filename": os.path.basename(file_path),
            "file_hash": file_hash,
            "url": f"http://127.0.0.1:8002/docs/{SOURCE_ID}/{os.path.basename(file_path)}",
            "source": SOURCE_ID
        }
        try:
            response = requests.post(API_URL, json=payload)
            print(f"[INFO] Sent chunk {idx}/{len(chunks)} → {response.status_code}")
            if response.status_code != 200:
                print(f"[ERROR] {response.text}")
        except Exception as e:
            print(f"[ERROR] Failed to send chunk {idx}: {e}")

def process_folder(folder_path: str):
    for path in Path(folder_path).rglob("*"):
        if path.is_file():
            print(f"[PROCESS] {path}")
            process_single_file(str(path))


# ───── Entry Point ───── #
if __name__ == "__main__":
    process_folder(FOLDER_TO_PROCESS)
